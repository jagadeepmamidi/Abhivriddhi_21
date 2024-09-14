from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
from cryptography.hazmat.backends import default_backend
from cryptography.fernet import Fernet
from typing import Union
from web3.datastructures import AttributeDict
import spacy
import streamlit as st
import os
import base64
import re
import json
import logging
import random
import string
from typing import List, Dict, Optional, Tuple
from datetime import datetime, timedelta
from io import StringIO
import sys
import hashlib
from web3 import Web3
import json
import streamlit as st
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import io
import pyrebase
import easyocr
from PIL import Image, ImageDraw



SENSITIVE_PATTERNS = {
        'email': r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}',
        'name': r'\b[A-Z][a-z]+\b',  # Simple pattern for names (first letter capitalized)
        'address': r'\d{1,5}\s\w+(\s\w+){1,3}',  # Simple pattern for addresses
        'card_number': r'\b\d{4}(-\d{4}){3}\b',  # Simple pattern for card numbers (XXXX-XXXX-XXXX-XXXX)
        'organization': r'\b[A-Z][a-zA-Z\s]+(?:Inc|LLC|Ltd|Corp|Co)\b',  # Pattern for organization names
    }

def mask_sensitive_data(image, patterns):
    reader = easyocr.Reader(['en'])
    results = reader.readtext(image)
    
    img = Image.open(io.BytesIO(image))
    draw = ImageDraw.Draw(img)
    
    # Check and mask sensitive data based on patterns
    for result in results:
        text = result[1]
        box = result[0]
        # Convert the box coordinates to integers
        box = [(int(x), int(y)) for (x, y) in box]
        for key, pattern in patterns.items():
            if re.search(pattern, text, re.IGNORECASE):
                draw.polygon(box, fill='black')
                break  # Masking once per detected box

    return img

def redact_personal_info(text, use_spacy=False, domain='general'):
    """
    Redacts sensitive information based on the chosen method (regex or spaCy) and domain.
    
    Args:
    text (str): The input text containing sensitive information.
    use_spacy (bool): Whether to use spaCy for redaction (if False, use regex).
    domain (str): The domain for which to apply specific redaction patterns.
    
    Returns:
    str: The redacted text.
    """
    
    # Domain-specific patterns
    domain_patterns = {
        'general': {
            'bank_account': r'\b\d{10,16}\b',
            'credit_card': r'\b(?:\d[ -]*?){13,16}\b',
            'tin': r'\b\d{3}-\d{2}-\d{4}\b',
            'salary_info': r'\b(\$|\₹)?\d{1,3}(,\d{3})*(\.\d{2})?\b',
            'ssn': r'\b\d{3}-\d{2}-\d{4}\b',
            'passport_number': r'\b[A-Z0-9]{6,9}\b',
            'driving_license': r'\b[A-Z0-9]{8,12}\b',
            'birthdate': r'\b(?:\d{1,2}[/-]\d{1,2}[/-]\d{2,4})\b',
            'national_id': r'\b\d{9,12}\b'
        },
        'financial': {
            'amount': r'\b(\$|£|€|₹)?\d{1,3}(?:,\d{3})*(?:\.\d{2})?\b',
            'credit_card': r'\b(?:\d[ -]*?){13,16}\b',
            'bank_account': r'\b\d{10,16}\b',
            'transaction_id': r'\b[Tt][Aa][Xx][Nn][Uu][Mm][Bb][Ee][Rr]-\d{8}\b'
        },
        'personal': {
            'phone_number': r'\b(?:\+?(\d{1,3}))?[-.●]?\(?(?:\d{1,4})\)?[-.●]?\d{1,4}[-.●]?\d{1,4}[-.●]?\d{1,9}\b',
            'email': r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
            'date': r'\b(?:\d{1,2}[\/\-]\d{1,2}[\/\-]\d{2,4}|\b\w+\s\d{1,2},?\s\d{4})\b',
            'username': r'\b(username|user|login):\s*[A-Za-z0-9._%+-]+\b',
            'password': r'\b(password|pwd|pass):\s*\S+\b'
        },
        'health': {
            'medical_record': r'\bMR-\d{6,8}\b',
            'patient_id': r'\bPID-\d{6,10}\b',
            'insurance_policy': r'\bINS-\d{8,12}\b'
        }
    }

    # Select patterns based on the domain
    patterns = domain_patterns.get(domain, domain_patterns['general'])

    if not use_spacy:
        # Regex-based redaction
        for label, pattern in patterns.items():
            text = re.sub(pattern, '[REDACTED]', text)

    else:
        # spaCy-based redaction
        doc = nlp(text)
        for ent in doc.ents:
            if ent.label_ in ["PERSON", "GPE", "ORG", "LAW", "LOC", "MONEY", "CARDINAL", "DATE", "TIME"]:
                text = text.replace(ent.text, '[REDACTED]')
    
    return text


def process_docx(file) -> str:
    doc = Document(io.BytesIO(file.read()))
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_docx(file):
    doc = Document(io.BytesIO(file))
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_pdf(file):
    pdf_document = fitz.open(stream=file, filetype="pdf")
    text = ""
    for page in pdf_document:
        text += page.get_text()
    return text

def extract_text_from_pptx(file):
    ppt = Presentation(io.BytesIO(file))
    text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def process_pptx(file) -> str:
    presentation = Presentation(io.BytesIO(file.read()))
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def process_pdf(file) -> str:
    # Open the PDF file
    pdf_document = fitz.open(stream=file.read(), filetype="pdf")
    
    # Extract text from each page
    text = ""
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        text += page.get_text()
    
    return text



logging.basicConfig(
    filename="data_protection_tool.log",
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
)

blockchain_url = (
    "http://127.0.0.1:7545"  # Ensure this matches your Ganache or node setup
)
web3 = Web3(Web3.HTTPProvider(blockchain_url))

contract_address = "0x41ed3031dA1C3f9E2560B2dAF7472cA862516D3D"  # Replace with your deployed contract address

# Use the user's suggested approach
compiled_contract_path = os.path.join(
    os.path.dirname(__file__), "../build/contracts/AuditLog.json"
)
with open(compiled_contract_path, "r") as file:
    audit_contract_json = json.load(file)
    audit_contract_abi = audit_contract_json["abi"]


contract = web3.eth.contract(address=contract_address, abi=audit_contract_abi)

# Set default account
web3.eth.default_account = web3.eth.accounts[0]

import os

def process_file(file):
    file_extension = os.path.splitext(file.name)[1].lower()
    if file_extension == ".txt":
        return file.read().decode("utf-8")
    elif file_extension == ".pdf":
        return process_pdf(file)
    elif file_extension == ".docx":
        return process_docx(file)
    elif file_extension == ".pptx":
        return process_pptx(file)
    else:
        raise ValueError("Unsupported file format")







def get_audit_logs():
    """
    Retrieve all audit logs from the blockchain.
    """
    try:
        logs = contract.functions.getLogs().call()
        for log in logs:
            print(
                f"User: {log[0]}, Hash: {log[1]}, Action: {log[2]}, Timestamp: {log[3]}"
            )
    except Exception as e:
        print(f"Error retrieving logs: {e}")


def serialize_web3_object(obj):
    if isinstance(obj, AttributeDict):
        return {k: serialize_web3_object(v) for k, v in obj.items()}
    if isinstance(obj, (bytes, bytearray)):
        return obj.hex()
    if isinstance(obj, (int, float, str, bool, type(None))):
        return obj
    return str(obj)

def add_audit_log(user_id: str, data_hash: str, action: str, timestamp: str, use_encryption: bool = False):
    try:
        # Include use_encryption in the log data
        tx_hash = contract.functions.addLog(user_id, data_hash, f"{action} (Encrypted: {use_encryption})", timestamp).transact()
        receipt = web3.eth.wait_for_transaction_receipt(tx_hash)

        print(f"Log added to blockchain. Transaction Hash: {tx_hash.hex()}")
        
        st.success(f"Log added to blockchain. Transaction Hash: {tx_hash.hex()}")
        
        # Display Transaction Receipt
        st.subheader("Transaction Receipt")
        
        # Create columns for better layout
        col1, col2 = st.columns(2)
        
        with col1:
            st.metric("Block Number", receipt.blockNumber)
            st.metric("Gas Used", receipt.gasUsed)
            st.metric("Status", "Success" if receipt.status == 1 else "Failed")
        
        with col2:
            st.metric("Transaction Index", receipt.transactionIndex)
            st.metric("Cumulative Gas Used", receipt.cumulativeGasUsed)
            st.metric("Effective Gas Price", receipt.effectiveGasPrice)
        
        # Display more detailed information in an expander
        with st.expander("Detailed Transaction Information"):
            st.json({
                "transactionHash": receipt.transactionHash.hex(),
                "blockHash": receipt.blockHash.hex(),
                "from": receipt["from"],
                "to": receipt.to,
                "contractAddress": receipt.contractAddress,
                "logs": [serialize_log(log) for log in receipt.logs],
                "logsBloom": receipt.logsBloom.hex()
            })

    except Exception as e:
        error_message = f"Error adding log to blockchain: {e}"
        print(error_message)
        st.error(error_message)
        raise

def serialize_log(log):
    return {
        "address": log.address,
        "topics": [topic.hex() for topic in log.topics],
        "data": log.data.hex(),
        "blockNumber": log.blockNumber,
        "transactionHash": log.transactionHash.hex(),
        "transactionIndex": log.transactionIndex,
        "blockHash": log.blockHash.hex(),
        "logIndex": log.logIndex
    }
    
def format_log_entry(entry):
    # Extract timestamp, method, and hash using regex
    timestamp_match = re.search(r"\d{4}-\d{2}-\d{2} \d{2}:\d{2}:\d{2},\d{3}", entry)
    method_match = re.search(r"Method: (Data Redaction|Data Masking)", entry)
    hash_match = re.search(r"Hash: (\w+)", entry)

    timestamp = timestamp_match.group(0) if timestamp_match else "Unknown"
    method = method_match.group(1) if method_match else "Unknown"
    hash_value = hash_match.group(1) if hash_match else "Unknown"

    return f"""
    <div style="margin-bottom: 10px; padding: 10px; border: 1px solid #ddd; border-radius: 5px;">
        <p style="margin: 0; font-weight: bold;">{timestamp}</p>
        <p style="margin: 0; color: #555;">Method: <span style="font-weight: bold;">{method}</span></p>
        <p style="margin: 0; color: #888;">Hash: <code>{hash_value}</code></p>
        <p style="margin: 0; color: #777;">{entry}</p>
    </div>
    """


def hash_data(text: str) -> str:
    return hashlib.sha256(text.encode()).hexdigest()


def log_protection_activity(
    user_id: str,
    original_text: str, 
    processed_text: str, 
    protection_method: str, 
    redaction_level: str = None, 
    use_encryption: bool = False,
    domain: str = None
):
    timestamp = datetime.now().isoformat()
    redacted_amount = len(original_text) - len(processed_text)
    data_hash = hash_data(processed_text)

    log_message = (
        f"{timestamp} | "
        f"User ID: {user_id} | "
        f"Method: {protection_method} | "
        f"Redacted: {redacted_amount} chars | "
        f"Hash: {data_hash} | "
        f"Redaction Level: {redaction_level if redaction_level else 'N/A'} | "
        f"Encryption Used: {'Yes' if use_encryption else 'No'} | "
        f"Domain: {domain if domain else 'N/A'}"
    )

    # Log to file
    logging.info(log_message)

    # Log to blockchain
    try:
        add_audit_log(user_id, data_hash, protection_method, timestamp, use_encryption)
        logging.info(f"Log added to blockchain successfully: {data_hash}")
    except Exception as e:
        logging.error(f"Error adding log to blockchain: {e}")


def process_text(user_id, text, protection_method, entity_types, custom_words, redaction_level=None, use_encryption=False, domain=None, use_spacy=False):
    original_text = text
    key = None
    encrypted_original = None
    
    if protection_method == "Data Redaction":
        result, modifications = redact_entities(text, entity_types, custom_words, redaction_level)
    elif protection_method == "Data Masking":
        result = mask_data(text, entity_types, custom_words)
    elif protection_method == "Data Anonymization":
        result = anonymize_data(text, entity_types, custom_words)
    elif protection_method == "Domain-Specific Redaction":
        result = redact_personal_info(text, use_spacy=use_spacy, domain=domain)
    else:
        raise ValueError("Invalid protection method")

    if use_encryption:
        key = generate_key()
        encrypted_original = encrypt_full_text(original_text, key)

    log_protection_activity(user_id, original_text, result, protection_method, redaction_level, use_encryption=use_encryption, domain=domain)

    return result, key, encrypted_original



DOWNLOAD_HISTORY_FILE = "download_history.json"
UPLOAD_FOLDER = "uploads"
FILE_RETENTION_DAYS = 30
DEFAULT_ENTITY_TYPES = ["PERSON", "ORG", "GPE", "DATE", "EMAIL"]
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
undo_stack = []
redo_stack = []


def add_email_matcher(nlp):
    email_pattern = r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b"
    matcher = spacy.matcher.Matcher(nlp.vocab)
    matcher.add("EMAIL", [[{"TEXT": {"REGEX": email_pattern}}]])
    return matcher


@st.cache_resource
def load_nlp_model():
    nlp = spacy.load("en_core_web_sm")
    email_matcher = add_email_matcher(nlp)
    return nlp, email_matcher


nlp, email_matcher = load_nlp_model()

def generate_key():
    return Fernet.generate_key()

def encrypt_full_text(text: str, key: bytes) -> str:
    f = Fernet(key)
    return f.encrypt(text.encode()).decode()

def decrypt_full_text(encrypted_text: str, key: bytes) -> str:
    f = Fernet(key)
    return f.decrypt(encrypted_text.encode()).decode()

def redact_entities(text: str, entity_types: List[str], custom_words: Optional[List[str]] = None, redaction_level: str = "Low") -> Tuple[str, List[Tuple[int, int, str, str]]]:
    doc = nlp(text)
    modifications = []

    def modify_entity(word: str) -> str:
        if redaction_level == "High":
            return "[Redacted]"
        elif redaction_level == "Medium":
            return f"[Redacted:{len(word)}]"
        else:
            return f"[Redacted:{word[:2]}...]"

    # Process custom words
    if custom_words:
        for word in custom_words:
            for match in re.finditer(re.escape(word), text, re.IGNORECASE):
                modified = modify_entity(word)
                modifications.append((match.start(), match.end(), modified, "CUSTOM"))

    # Process emails
    if "EMAIL" in entity_types:
        email_matches = email_matcher(doc)
        for _, start, end in email_matches:
            modified = modify_entity(text[start:end])
            modifications.append((start, end, modified, "EMAIL"))

    # Process other entities
    for ent in doc.ents:
        if ent.label_ in entity_types:
            modified = modify_entity(ent.text)
            modifications.append((ent.start_char, ent.end_char, modified, ent.label_))

    # Apply modifications in reverse order to preserve indices
    modifications.sort(key=lambda x: x[1], reverse=True)
    for start, end, replacement, label in modifications:
        text = text[:start] + replacement + text[end:]

    return text, modifications


def generate_fake_data(entity_type: str) -> str:
    if entity_type == "PERSON":
        return "John Doe"
    elif entity_type == "ORG":
        return "ACME Corporation"
    elif entity_type == "GPE":
        return "Anytown"
    elif entity_type == "DATE":
        return "01/01/2000"
    elif entity_type == "EMAIL":
        return f"user{random.randint(1000, 9999)}@example.com"
    else:
        return "".join(random.choices(string.ascii_letters + string.digits, k=10))


def mask_data(
    text: str, entity_types: List[str], custom_words: Optional[List[str]] = None
) -> str:
    doc = nlp(text)
    masked_text = text

    if custom_words:
        for word in custom_words:
            masked_text = re.sub(
                re.escape(word),
                generate_fake_data("CUSTOM"),
                masked_text,
                flags=re.IGNORECASE,
            )

    if "EMAIL" in entity_types:
        email_matches = email_matcher(doc)
        for _, start, end in reversed(email_matches):
            fake_email = generate_fake_data("EMAIL")
            masked_text = masked_text[:start] + fake_email + masked_text[end:]

    for ent in reversed(doc.ents):
        if ent.label_ in entity_types:
            fake_data = generate_fake_data(ent.label_)
            masked_text = (
                masked_text[: ent.start_char] + fake_data + masked_text[ent.end_char :]
            )

    return masked_text


def anonymize_data(
    text: str, entity_types: List[str], custom_words: Optional[List[str]] = None
) -> str:
    doc = nlp(text)
    anonymized_text = text

    def generate_anonymous_id() -> str:
        return "".join(random.choices(string.ascii_uppercase + string.digits, k=8))

    if custom_words:
        for word in custom_words:
            anonymized_text = re.sub(
                re.escape(word),
                generate_anonymous_id(),
                anonymized_text,
                flags=re.IGNORECASE,
            )

    if "EMAIL" in entity_types:
        email_matches = email_matcher(doc)
        for _, start, end in reversed(email_matches):
            anonymous_id = generate_anonymous_id()
            anonymized_text = (
                anonymized_text[:start]
                + anonymous_id
                + "@anon.com"
                + anonymized_text[end:]
            )

    for ent in reversed(doc.ents):
        if ent.label_ in entity_types:
            anonymous_id = generate_anonymous_id()
            anonymized_text = (
                anonymized_text[: ent.start_char]
                + anonymous_id
                + anonymized_text[ent.end_char :]
            )

    return anonymized_text


def get_entity_counts(text: str) -> Dict[str, int]:
    doc = nlp(text)
    entity_counts = {}
    for ent in doc.ents:
        entity_counts[ent.label_] = entity_counts.get(ent.label_, 0) + 1

    email_matches = email_matcher(doc)
    entity_counts["EMAIL"] = len(email_matches)

    return entity_counts


def get_download_link(text, file_name, link_text="Download"):
    # Create a BytesIO object to hold the text data
    buffer = io.BytesIO()
    
    # Write text to the BytesIO buffer directly
    buffer.write(text.encode('utf-8'))
    
    # Seek to the beginning of the buffer to read its contents
    buffer.seek(0)
    
    # Encode the buffer contents to base64
    b64 = base64.b64encode(buffer.read()).decode()
    
    # Create a download link
    href = f'<a href="data:file/txt;base64,{b64}" download="{file_name}">{link_text}</a>'
    return href



def save_download_history(filename: str):
    try:
        history = []
        if os.path.exists(DOWNLOAD_HISTORY_FILE):
            with open(DOWNLOAD_HISTORY_FILE, "r") as f:
                history = json.load(f)

        history.append({"filename": filename, "timestamp": datetime.now().isoformat()})

        with open(DOWNLOAD_HISTORY_FILE, "w") as f:
            json.dump(history, f)
        logging.info(f"Download history updated: {filename}")
    except Exception as e:
        logging.error(f"Error saving download history: {e}")
        st.error(f"Error saving download history: {e}")


def get_download_history() -> List[Dict[str, str]]:
    if os.path.exists(DOWNLOAD_HISTORY_FILE):
        try:
            with open(DOWNLOAD_HISTORY_FILE, "r") as f:
                return json.load(f)
        except Exception as e:
            logging.error(f"Error retrieving download history: {e}")
            st.error(f"Error retrieving download history: {e}")
    return []


def cleanup_old_files():
    try:
        current_time = datetime.now()
        for filename in os.listdir(UPLOAD_FOLDER):
            file_path = os.path.join(UPLOAD_FOLDER, filename)
            file_modified = datetime.fromtimestamp(os.path.getmtime(file_path))
            if current_time - file_modified > timedelta(days=FILE_RETENTION_DAYS):
                os.remove(file_path)
                logging.info(f"Deleted old file: {filename}")
    except Exception as e:
        logging.error(f"Error during file cleanup: {e}")

firebaseConfig = {
  "apiKey": "AIzaSyBSxQrlZX9-1Hwkse8JF3RNKFojt519wQs",
  "authDomain": "pachack-49b6f.firebaseapp.com",
  "databaseURL": "https://pachack-49b6f-default-rtdb.asia-southeast1.firebasedatabase.app",
  "projectId": "pachack-49b6f",
  "storageBucket": "pachack-49b6f.appspot.com",
  "messagingSenderId": "875463536726",
  "appId": "1:875463536726:web:d5d1a2a2e39720c00a4bf7",
  "measurementId": "G-YQZV45GGGF"
}

# Initialize Firebase
firebase = pyrebase.initialize_app(firebaseConfig)
auth = firebase.auth()
db = firebase.database()

# Streamlit app


# Sidebar for Login or Signup




def main():
    
    
    st.title("RE-DACT")
    st.text("Enhanced Data Protection Tool")
    st.text("Built with Streamlit and spaCy")

    if 'user' not in st.session_state:
        st.session_state['user'] = None

    if not st.session_state['user']:
        choice = st.sidebar.selectbox("Login/Signup", ["Login", "Signup"])
        
        if choice == "Signup":
            with st.form("signup_form"):
                st.subheader("Signup")
                username = st.text_input("Enter username")
                email = st.text_input("Enter email")
                password = st.text_input("Enter password", type="password")
                signup_button = st.form_submit_button("Signup")

                if signup_button:
                    try:
                        user = auth.create_user_with_email_and_password(email, password)
                        st.success("Account created successfully!")
                        user_id = user['localId']
                        db.child("users").child(user_id).set({
                            "username": username,
                            "email": email
                        })
                        st.session_state['user'] = user
                    except Exception as e:
                        st.error(f"Error creating account: {e}")

        elif choice == "Login":
            with st.form("login_form"):
                st.subheader("Login")
                email = st.text_input("Enter email")
                password = st.text_input("Enter password", type="password")
                login_button = st.form_submit_button("Login")

                if login_button:
                    try:
                        user = auth.sign_in_with_email_and_password(email, password)
                        st.success("Login successful!")
                        st.session_state['user'] = user
                        user_id = user['localId']
                        user_data = db.child("users").child(user_id).get().val()
                        if user_data:
                            st.write(f"Welcome, {user_data.get('username')}!")
                        else:
                            st.write("No additional user data found.")
                    except Exception as e:
                        st.error(f"Error logging in: {e}")

    if st.session_state['user']:
        
        
        
        activities = [
            "Data Protection",
            "Entity Analysis",
            "Downloads",
            "About",
            "View Logs",
            "Decrypt Text",
        ]
        choice = st.sidebar.selectbox("Select Task", activities)

        if choice == "Data Protection":
            image_path = ".//one.jpg"  # Update this path to the correct location of your image
            if os.path.exists(image_path):
                try:
                    img = Image.open(image_path)
                    st.image(img, use_column_width=True)
                except Exception as e:
                    st.error(f"Error loading image: {e}")
            else:
                st.warning("Image file not found. Continuing without the image.")   
    
        if choice == "Data Protection":
            image_path = ".//two.jpg"  # Update this path to the correct location of your image
            if os.path.exists(image_path):
                try:
                    img = Image.open(image_path)
                    st.image(img, use_column_width=True)
                except Exception as e:
                    st.error(f"Error loading image: {e}")
            else:
                st.warning("Image file not found. Continuing without the image.")   
        
        if choice == "Data Protection":
            image_path = ".//three.jpg"  # Update this path to the correct location of your image
            if os.path.exists(image_path):
                try:
                    img = Image.open(image_path)
                    st.image(img, use_column_width=True)
                except Exception as e:
                    st.error(f"Error loading image: {e}")
            else:
                st.warning("Image file not found. Continuing without the image.") 
                
        if choice == "Data Protection":
            image_path = ".//four.jpg"  # Update this path to the correct location of your image
            if os.path.exists(image_path):
                try:
                    img = Image.open(image_path)
                    st.image(img, use_column_width=True)
                except Exception as e:
                    st.error(f"Error loading image: {e}")
            else:
                st.warning("Image file not found. Continuing without the image.")    
            
            
            
            st.subheader("Data Protection Options")

            
            
            uploaded_file = st.file_uploader("Choose a file", type=["txt", "csv", "pdf", "docx", "pptx"])
            if uploaded_file is not None:
                file_type = uploaded_file.type

                if file_type == "application/pdf":
                    rawtext = extract_text_from_pdf(uploaded_file.read())
                elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    rawtext = extract_text_from_docx(uploaded_file.read())
                elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    rawtext = extract_text_from_pptx(uploaded_file.read())
                elif file_type == "text/plain" or file_type == "text/csv":
                    rawtext = uploaded_file.getvalue().decode("utf-8")
                else:
                    st.error("Unsupported file type")
                    rawtext = ""

                if rawtext:
                    st.text_area("File content", rawtext, height=300)
            else:
                rawtext = st.text_area("Or enter text", "Type Here", height=300)

            all_entity_types = list(nlp.get_pipe("ner").labels) + ["EMAIL"]
            entity_types = st.multiselect(
                "Select entity types to protect",
                all_entity_types,
                default=DEFAULT_ENTITY_TYPES,
            )

            custom_words = st.text_input("Enter custom words to protect (comma-separated)")
            custom_word_list = (
                [word.strip() for word in custom_words.split(",") if word.strip()]
                if custom_words
                else None
            )

            st.title("Sensitive Data Masking")

            uploaded_file = st.file_uploader("Choose an image...", type=["jpg", "jpeg", "png"])

            if uploaded_file is not None:
                # Load the uploaded image file
                image = uploaded_file.read()

                # Mask sensitive data in the image
                masked_image = mask_sensitive_data(image, SENSITIVE_PATTERNS)

                # Display the masked image
                st.image(masked_image, caption='Processed Image', use_column_width=True)

                # Convert the masked image to a format suitable for download
                img_byte_arr = io.BytesIO()
                masked_image.save(img_byte_arr, format='PNG')
                img_byte_arr.seek(0)

                # Create a download button for the redacted image
                st.download_button(
                    label="Download Redacted Image",
                    data=img_byte_arr,
                    file_name="redacted_image.png",
                    mime="image/png"
                )
            
            protection_method = st.radio(
        "Select Protection Method",
        ("Data Redaction", "Data Masking", "Data Anonymization", "Domain-Specific Redaction"),
    )
    
            redaction_level = None
            use_encryption = False
            domain = None
            use_spacy = False

            if protection_method == "Data Redaction":
                redaction_level = st.radio(
                    "Select Redaction Level", ("High", "Medium", "Low")
                )
            elif protection_method == "Domain-Specific Redaction":
                domain = st.selectbox("Select the domain for redaction:", options=['general', 'financial', 'personal', 'health'])
                use_spacy = st.checkbox("Use spaCy for redaction (otherwise regex will be used)", value=False)
            
            use_encryption = st.checkbox("Use encryption", value=False)

            if st.button("Process"):
                if not rawtext or rawtext == "Type Here":
                    st.error("Please enter some text to protect or upload a file.")
                elif not entity_types and not custom_word_list:
                    st.error("Please select at least one entity type to protect or enter custom words.")
                else:
                    with st.spinner("Processing text..."):
                        user_id = st.session_state['user']['localId']
                        result, key, encrypted_original = process_text(
                            user_id,
                            rawtext, 
                            protection_method, 
                            entity_types, 
                            custom_word_list, 
                            redaction_level,
                            use_encryption,
                            domain,
                            use_spacy
                        )

                    # Log to blockchain
                    data_hash = hash_data(result)
                    timestamp = datetime.now().isoformat()
                    add_audit_log(user_id, data_hash, protection_method, timestamp, use_encryption)
                    
                    st.write("Processed Text:")
                    st.write(result)

                    if use_encryption and key and encrypted_original:
                        st.write("Encryption Key (save this to decrypt later):")
                        st.code(key)

                        st.write("Encrypted Original (store this securely):")
                        st.code(encrypted_original)

                    filename = f"{protection_method.lower().replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
                    st.markdown(
                        get_download_link(result, filename, f"Download {protection_method} Text"),
                        unsafe_allow_html=True,
                    )

                    save_download_history(filename)

        elif choice == "Decrypt Text":
            st.subheader("Decrypt Text")
            encrypted_text = st.text_area("Enter text with encrypted parts", height=200)
            encryption_key = st.text_input("Enter encryption key")

            if st.button("Decrypt"):
                try:
                    decrypted_text = decrypt_full_text(encrypted_text, encryption_key.encode())
                    st.write("Decrypted Text:")
                    st.text_area("Result", decrypted_text, height=300)
                except Exception as e:
                    st.error(f"Decryption failed: {str(e)}")

        elif choice == "Entity Analysis":
            st.subheader("Entity Analysis")
            rawtext = st.text_area("Enter text for analysis", "Type Here", height=300)
            if st.button("Analyze"):
                if not rawtext or rawtext == "Type Here":
                    st.error("Please enter some text to analyze.")
                else:
                    with st.spinner("Analyzing text..."):
                        entity_counts = get_entity_counts(rawtext)
                    st.write("Entity Counts:")
                    for entity, count in entity_counts.items():
                        st.write(f"{entity}: {count}")

        elif choice == "Downloads":
            st.subheader("Download History")
            history = get_download_history()
            if history:
                for item in history:
                    st.write(f"{item['timestamp']}: {item['filename']}")
                    file_path = os.path.join(UPLOAD_FOLDER, item["filename"])
                    if os.path.exists(file_path):
                        with open(file_path, "r") as f:
                            st.download_button(
                                f"Download {item['filename']}", f.read(), item["filename"]
                            )
                    else:
                        st.write(f"File {item['filename']} no longer exists.")
            else:
                st.write("No download history available.")

        # elif choice == "View Logs":
        #   st.subheader("View Logs")

        # Path to your log file
        #  log_file_path = ".//data_protection_tool.log"

        # if os.path.exists(log_file_path):
        #    with open(log_file_path, 'r') as file:
        #       log_entries = file.readlines()

        # Display log entries with custom formatting
        #  for entry in log_entries:
        #     formatted_entry = format_log_entry(entry.strip())
        #    st.markdown(formatted_entry, unsafe_allow_html=True)
        # else:
        #   st.error("Log file not found.")

        elif choice == "View Logs":
            st.subheader("View Blockchain Logs")

            if st.button("Fetch Blockchain Logs"):
                try:
                    user_id = st.session_state['user']['localId']
                    user_logs = contract.functions.getUserLogs(user_id).call()

                    if user_logs:
                        for log in user_logs:
                            st.markdown(
                            f"""
                            ---
                            **Hash**: {log[1]}  
                            **Action**: {log[2]}  
                            **Timestamp**: {log[3]}
                            """
                        )
                    else:
                        st.write("No logs found for this user.")
                except Exception as e:
                    st.error(f"Error fetching logs: {e}")
                    print(f"Error details: {str(e)}")
            
        elif choice == "About":
            st.subheader("About")
            st.write(
                "This is an enhanced data protection tool built with Streamlit and spaCy."
            )
            st.write("Features include:")
            st.write(
                "- Data Redaction: Obscures or blacks out sensitive or confidential information in a document"
            )
            st.write(
                "- Data Masking: Replaces authentic information with fake one, but with the same structure"
            )
            st.write(
                "- Data Anonymization: Erases/encrypts identifiers in a document so identification is not possible"
            )
            st.write("- Entity analysis")
            st.write("- Download history and re-download of previous files")
            st.write("- File upload support")
            st.write("- Email detection and protection")

    # Run cleanup job
    cleanup_old_files()         
if __name__ == "__main__":
    main()
